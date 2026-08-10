"""技能过滤与 MCP 工具管理测试"""
import json
from types import SimpleNamespace

import pytest

import core.generator as cg
import core.skills as sk
from config import settings
from core.generator import Generator
from core.mcp_manager import McpManager, _field


@pytest.fixture
def gen(monkeypatch):
    monkeypatch.setattr(settings, "deepseek_api_key", "test-key")
    return Generator()


# ── 技能过滤 ──


def test_build_tools_filters_by_enabled_names():
    tools = sk.build_tools(["calculate", "get_current_time"], None)
    names = [t["function"]["name"] for t in tools]
    assert set(names) == {"calculate", "get_current_time"}


def test_build_tools_empty_disables_all_builtins():
    assert sk.build_tools([], None) == []


def test_build_tools_none_keeps_all_builtins():
    assert len(sk.build_tools(None, None)) == len(sk.TOOLS)


def test_is_builtin_tool():
    assert sk.is_builtin_tool("calculate")
    assert not sk.is_builtin_tool("mcp__whatever")


# ── MCP 管理器 ──


def test_mcp_field_dict_and_object():
    assert _field({"name": "a"}, "name") == "a"
    assert _field(SimpleNamespace(name="b"), "name") == "b"
    assert _field({}, "missing", "x") == "x"


def test_mcp_manager_call_unknown_tool_returns_error():
    m = McpManager()
    m.configure([{"name": "srv", "url": "http://x"}])
    out = m.call_tool("srv__nope", {})
    assert "未知 MCP 工具" in out


def test_mcp_manager_configure_dict_and_object():
    m = McpManager()
    m.configure([{"name": "a", "url": "http://a"}, SimpleNamespace(name="b", url="http://b")])
    assert m._server_configs == {"a": "http://a", "b": "http://b"}


# ── 生成器按配置构建工具并路由 MCP 调用 ──


def test_generate_uses_filtered_tools(gen):
    calls = {}

    def create(**k):
        calls.update(k)
        return SimpleNamespace(choices=[SimpleNamespace(
            message=SimpleNamespace(content="ok", tool_calls=None))])

    gen.client = SimpleNamespace(
        chat=SimpleNamespace(completions=SimpleNamespace(create=create)))

    answer, _, _ = gen.generate("你好", [], skill_names=["calculate"])

    assert answer == "ok"
    names = [t["function"]["name"] for t in calls["tools"]]
    assert names == ["calculate"]


def test_generate_dispatches_mcp_tool(gen, monkeypatch):
    fake_mcp = SimpleNamespace(
        list_tools=lambda servers: [{
            "type": "function",
            "function": {
                "name": "srv__weather",
                "description": "查询天气",
                "parameters": {"type": "object", "properties": {}},
            },
        }],
        call_tool=lambda name, args: "MCP 天气结果：晴")
    monkeypatch.setattr("core.mcp_manager.mcp_manager", fake_mcp)

    creates = []

    def create(**k):
        creates.append(k)
        if len(creates) == 1:
            tc = SimpleNamespace(
                id="c1",
                function=SimpleNamespace(
                    name="srv__weather",
                    arguments=json.dumps({"city": "上海"})))
            return SimpleNamespace(choices=[SimpleNamespace(
                message=SimpleNamespace(content=None, tool_calls=[tc]))])
        return SimpleNamespace(choices=[SimpleNamespace(
            message=SimpleNamespace(content="最终回答", tool_calls=None))])

    gen.client = SimpleNamespace(
        chat=SimpleNamespace(completions=SimpleNamespace(create=create)))

    answer, _, _ = gen.generate(
        "上海天气怎么样", [], mcp_servers=[{"name": "srv", "url": "http://x"}])

    assert answer == "最终回答"
    tool_msgs = [m for m in creates[1]["messages"] if m.get("role") == "tool"]
    assert any("MCP 天气结果" in m["content"] for m in tool_msgs)


# ── 新增技能（电商/爬虫/报表/汇总） ──


def test_tools_include_new_skills():
    names = {t["function"]["name"] for t in sk.TOOLS}
    for n in ["barcode_lookup", "exchange_convert", "github_search", "arxiv_search",
              "hn_search", "pypi_info", "stock_quote", "stock_digest", "make_table"]:
        assert n in names


def test_make_table_markdown():
    out = sk.execute_tool("make_table", {
        "data": '[{"名称":"A","销量":100},{"名称":"B","销量":200}]'})
    assert "| 名称 | 销量 |" in out
    assert "| A | 100 |" in out


def test_make_table_csv():
    out = sk.execute_tool("make_table", {"data": "[[1,2],[3,4]]", "format": "csv"})
    assert out.splitlines()[0] == "列1,列2"
    assert "1,2" in out


def test_make_table_invalid_json():
    out = sk.execute_tool("make_table", {"data": "not-json"})
    assert "JSON" in out


def test_barcode_lookup_invalid_input_no_network():
    out = sk.execute_tool("barcode_lookup", {"barcode": "abc"})
    assert "8-14" in out


def test_pypi_invalid_name():
    out = sk.execute_tool("pypi_info", {"name": "../etc"})
    assert "不合法" in out


def test_stock_invalid_symbol():
    out = sk.execute_tool("stock_quote", {"symbol": "bad symbol!"})
    assert "失败" in out or "无效" in out


def test_exchange_convert_invalid_amount():
    out = sk.execute_tool("exchange_convert", {"from": "USD", "to": "CNY", "amount": "abc"})
    assert "金额" in out


# ── web_extract 动态爬虫技能 ──


def test_tools_include_web_extract():
    names = {t["function"]["name"] for t in sk.TOOLS}
    assert "web_extract" in names


def test_web_extract_disabled_without_allowlist(monkeypatch):
    monkeypatch.setattr(settings, "crawl_allowlist", [])
    monkeypatch.setattr(settings, "url_fetch_allowlist", [])
    out = sk.execute_tool("web_extract", {"url": "https://example.com/"})
    assert "未启用" in out


def test_web_extract_rejects_internal_ip(monkeypatch):
    monkeypatch.setattr(settings, "crawl_allowlist", ["example.com"])
    monkeypatch.setattr("core.skills.socket.getaddrinfo",
                        lambda h, p: [(2, 1, 6, "", ("127.0.0.1", 0))])
    out = sk.execute_tool("web_extract", {"url": "https://example.com/"})
    assert "内网" in out


def test_web_extract_extracts_markdown(monkeypatch):
    monkeypatch.setattr(settings, "crawl_allowlist", ["example.com"])
    monkeypatch.setattr("core.skills.socket.getaddrinfo",
                        lambda h, p: [(2, 1, 6, "", ("93.184.216.34", 0))])
    fake = SimpleNamespace(
        text="<html><body><article><h1>Title</h1><p>Hello world content.</p></article></body></html>",
        status_code=200, raise_for_status=lambda: None, is_redirect=False, headers={})
    monkeypatch.setattr(sk.requests, "get", lambda *a, **k: fake)
    monkeypatch.setattr("trafilatura.extract",
                        lambda html, **k: "## Title\n\nHello world content.")
    out = sk.execute_tool("web_extract", {"url": "https://example.com/"})
    assert "Hello world content" in out


# ── 文档提取 / 数据处理 / 工具类技能 ──


def test_tools_include_document_and_util_skills():
    names = {t["function"]["name"] for t in sk.TOOLS}
    for n in ["docx_extract", "xlsx_extract", "pptx_extract", "pdf_extract",
              "csv_tools", "text_stats", "ip_lookup", "mermaid_chart",
              "qr_generate", "today_hot"]:
        assert n in names


def test_text_stats():
    out = sk.execute_tool("text_stats", {"text": "你好 world 123。第二句！"})
    assert "中文字 5" in out and "数字 3" in out and "句子 2" in out


def test_csv_tools_table_and_summary_and_dedupe():
    csv_text = "名称,销量\nA,100\nB,200\nA,100\n"
    t = sk.execute_tool("csv_tools", {"csv_text": csv_text, "op": "table"})
    assert "| 名称 | 销量 |" in t and "| A | 100 |" in t
    s = sk.execute_tool("csv_tools", {"csv_text": csv_text, "op": "summary"})
    assert "合计 400" in s and "平均 133.33" in s
    d = sk.execute_tool("csv_tools", {"csv_text": csv_text, "op": "dedupe"})
    assert d.count("| A | 100 |") == 1


def test_mermaid_chart_pie_and_flowchart():
    pie = sk.execute_tool("mermaid_chart",
                          {"chart_type": "pie", "data": '{"苹果":3,"香蕉":2}'})
    assert pie.startswith("```mermaid") and '"苹果" : 3' in pie
    flow = sk.execute_tool("mermaid_chart",
                           {"chart_type": "flowchart",
                            "data": '{"A":"开始","B":"结束","edges":[["A","B"]]}'})
    assert "A --> B" in flow


def test_docx_xlsx_pptx_pdf_extract(tmp_path, monkeypatch):
    monkeypatch.setattr(settings, "file_access_dirs", [str(tmp_path)])
    # Word
    from docx import Document
    d = Document()
    d.add_paragraph("这是Word测试段落")
    d.save(str(tmp_path / "t.docx"))
    out = sk.execute_tool("docx_extract", {"path": str(tmp_path / "t.docx")})
    assert "这是Word测试段落" in out
    # Excel
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["名称", "销量"])
    ws.append(["A", 100])
    ws.append(["B", 200])
    wb.save(str(tmp_path / "t.xlsx"))
    out = sk.execute_tool("xlsx_extract", {"path": str(tmp_path / "t.xlsx"), "op": "summary"})
    assert "合计 300" in out
    # PPT
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(0, 0, 4000000, 400000)
    tb.text = "这是PPT测试文本"
    prs.save(str(tmp_path / "t.pptx"))
    out = sk.execute_tool("pptx_extract", {"path": str(tmp_path / "t.pptx")})
    assert "这是PPT测试文本" in out
    # PDF
    import pymupdf
    pdf = pymupdf.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Hello PDF content")
    pdf.save(str(tmp_path / "t.pdf"))
    pdf.close()
    out = sk.execute_tool("pdf_extract", {"path": str(tmp_path / "t.pdf")})
    assert "Hello PDF content" in out


def test_ip_lookup_invalid_ip():
    out = sk.execute_tool("ip_lookup", {"ip": "999.1.1.1"})
    assert "不合法" in out


def test_qr_generate_writes_svg(tmp_path, monkeypatch):
    monkeypatch.setattr(settings, "icon_output_dir", str(tmp_path))
    monkeypatch.setattr(settings, "chart_base_url", "http://localhost:8080")
    out = sk.execute_tool("qr_generate", {"content": "https://example.com"})
    assert "qr_" in out and out.startswith("<img")
    assert list(tmp_path.glob("qr_*.svg"))


def test_today_hot_mocked(monkeypatch):
    fake = SimpleNamespace(
        json=lambda: {"data": [
            {"title": "热点新闻A", "hot": "100万", "url": "https://x"},
            {"title": "热点新闻B", "hot": "80万", "url": "https://y"},
        ]},
        raise_for_status=lambda: None)
    monkeypatch.setattr(sk.requests, "get", lambda *a, **k: fake)
    out = sk.execute_tool("today_hot", {"topic": "weibo"})
    assert "热点新闻A" in out and "微博热榜" in out
