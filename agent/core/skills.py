"""Skill 工具集 —— LLM 可调用的外部能力（全部基于免费 API）

所有工具均通过 OpenAI function calling 协议注册，LLM 自动选择调用。
收费情况：除 DeepSeek 调用费用外，所有工具 API 均免费。
"""
import ast
import ipaddress
import logging
import requests
import math
import re
import secrets
import socket
import string
from datetime import datetime, timedelta
from pathlib import Path
from typing import Dict, Any
from urllib.parse import urljoin, urlparse

from config import settings

logger = logging.getLogger(__name__)

# ── 工具定义 ──
TOOLS = [
  {
    "type": "function",
    "function": {
      "name": "get_current_time",
      "description": "获取当前日期和时间",
      "parameters": {
        "type": "object",
        "properties": {},
        "required": []
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "calculate",
      "description": "执行数学计算",
      "parameters": {
        "type": "object",
        "properties": {
          "expression": {
            "type": "string"
          }
        },
        "required": [
          "expression"
        ]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "web_search",
      "description": "搜索互联网",
      "parameters": {
        "type": "object",
        "properties": {
          "query": {
            "type": "string"
          }
        },
        "required": [
          "query"
        ]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "url_fetch",
      "description": "抓取网页内容",
      "parameters": {
        "type": "object",
        "properties": {
          "url": {
            "type": "string"
          }
        },
        "required": [
          "url"
        ]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "get_weather",
      "description": "查询天气",
      "parameters": {
        "type": "object",
        "properties": {
          "city": {
            "type": "string"
          }
        },
        "required": [
          "city"
        ]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "get_exchange_rate",
      "description": "查询汇率",
      "parameters": {
        "type": "object",
        "properties": {
          "from_currency": {
            "type": "string"
          },
          "to_currency": {
            "type": "string"
          }
        },
        "required": [
          "from_currency",
          "to_currency"
        ]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "wikipedia_lookup",
      "description": "查询百科",
      "parameters": {
        "type": "object",
        "properties": {
          "query": {
            "type": "string"
          },
          "lang": {
            "type": "string"
          }
        },
        "required": [
          "query"
        ]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "date_calc",
      "description": "日期计算",
      "parameters": {
        "type": "object",
        "properties": {
          "date1": {
            "type": "string"
          },
          "date2": {
            "type": "string"
          },
          "add_days": {
            "type": "integer"
          }
        },
        "required": [
          "date1"
        ]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "generate_password",
      "description": "生成密码",
      "parameters": {
        "type": "object",
        "properties": {
          "length": {
            "type": "integer"
          }
        },
        "required": []
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "make_chart",
      "description": "当用户要求画图/做图表/可视化数据时必须调用。生成柱状图(bar)/横向柱状图(barh)/折线图(line)/饼图(pie)，返回可直接嵌入的图片HTML标签",
      "parameters": {
        "type": "object",
        "properties": {
          "type": {
            "type": "string",
            "description": "图表类型",
            "enum": ["bar", "barh", "line", "pie"],
          },
          "labels": {
            "type": "string",
            "description": "标签，逗号分隔，如：一月,二月,三月"
          },
          "data": {
            "type": "string",
            "description": "数据值，逗号分隔，如：100,200,150"
          },
          "title": {
            "type": "string",
            "description": "图表标题（可选）"
          },
          "unit": {
            "type": "string",
            "description": "数据单位（可选），如：万元,%,人"
          }
        },
        "required": [
          "type",
          "labels",
          "data"
        ]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "news_headlines",
      "description": "获取新闻头条",
      "parameters": {
        "type": "object",
        "properties": {
          "topic": {
            "type": "string"
          }
        },
        "required": []
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "make_icon",
      "description": "生成SVG图标",
      "parameters": {
        "type": "object",
        "properties": {
          "description": {
            "type": "string"
          },
          "size": {
            "type": "integer"
          }
        },
        "required": [
          "description"
        ]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "barcode_lookup",
      "description": "查询商品条码信息（名称/品牌/规格/营养等级），数据来自 Open Food Facts",
      "parameters": {
        "type": "object",
        "properties": {
          "barcode": {"type": "string", "description": "8-14 位商品条码数字"}
        },
        "required": ["barcode"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "exchange_convert",
      "description": "汇率换算：把金额从一种货币换算成另一种货币",
      "parameters": {
        "type": "object",
        "properties": {
          "from": {"type": "string", "description": "源货币代码，如 USD"},
          "to": {"type": "string", "description": "目标货币代码，如 CNY"},
          "amount": {"type": "number", "description": "金额"}
        },
        "required": ["from", "to", "amount"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "github_search",
      "description": "搜索 GitHub 仓库（按星标排序）",
      "parameters": {
        "type": "object",
        "properties": {
          "query": {"type": "string", "description": "搜索关键词"}
        },
        "required": ["query"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "arxiv_search",
      "description": "检索 arXiv 学术论文（返回标题与摘要）",
      "parameters": {
        "type": "object",
        "properties": {
          "query": {"type": "string", "description": "论文关键词"}
        },
        "required": ["query"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "hn_search",
      "description": "检索 Hacker News 热门讨论",
      "parameters": {
        "type": "object",
        "properties": {
          "query": {"type": "string", "description": "搜索关键词"}
        },
        "required": ["query"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "pypi_info",
      "description": "查询 PyPI Python 包信息（版本/描述/主页）",
      "parameters": {
        "type": "object",
        "properties": {
          "name": {"type": "string", "description": "包名"}
        },
        "required": ["name"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "stock_quote",
      "description": "查询单只股票/指数实时行情（Stooq 免费源，如 aapl.us / ^spx）",
      "parameters": {
        "type": "object",
        "properties": {
          "symbol": {"type": "string", "description": "股票代码，如 aapl.us"}
        },
        "required": ["symbol"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "stock_digest",
      "description": "查询多只股票行情并汇总成 Markdown 报表",
      "parameters": {
        "type": "object",
        "properties": {
          "symbols": {"type": "string", "description": "逗号分隔的股票代码，如 aapl.us,msft.us"}
        },
        "required": ["symbols"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "make_table",
      "description": "把 JSON 数据转换成 Markdown 表格或 CSV（本地生成，无需网络）",
      "parameters": {
        "type": "object",
        "properties": {
          "data": {"type": "string", "description": "JSON 数组，如 [{\"名称\":\"A\",\"销量\":100}]"},
          "format": {"type": "string", "description": "md 或 csv，默认 md"}
        },
        "required": ["data"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "web_extract",
      "description": "抓取并提取网页正文为 Markdown（支持表格，可选保留链接；需域名在白名单）",
      "parameters": {
        "type": "object",
        "properties": {
          "url": {"type": "string", "description": "网页地址"},
          "extract_links": {"type": "boolean", "description": "是否保留链接（默认 false）"}
        },
        "required": ["url"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "docx_extract",
      "description": "提取 Word(.docx) 文档文本与表格（路径需在允许目录内）",
      "parameters": {
        "type": "object",
        "properties": {
          "path": {"type": "string", "description": "docx 文件路径"}
        },
        "required": ["path"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "xlsx_extract",
      "description": "读取 Excel(.xlsx)：表格预览或数值汇总（合计/平均/最大/最小）",
      "parameters": {
        "type": "object",
        "properties": {
          "path": {"type": "string", "description": "xlsx 文件路径"},
          "op": {"type": "string", "description": "table 或 summary，默认 table"}
        },
        "required": ["path"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "pptx_extract",
      "description": "提取 PowerPoint(.pptx) 每页文本",
      "parameters": {
        "type": "object",
        "properties": {
          "path": {"type": "string", "description": "pptx 文件路径"}
        },
        "required": ["path"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "pdf_extract",
      "description": "提取 PDF 文本与表格（前 20 页）",
      "parameters": {
        "type": "object",
        "properties": {
          "path": {"type": "string", "description": "pdf 文件路径"}
        },
        "required": ["path"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "csv_tools",
      "description": "CSV 数据处理：表格预览 / 数值汇总 / 去重（传入 CSV 内容或文件路径）",
      "parameters": {
        "type": "object",
        "properties": {
          "csv_text": {"type": "string", "description": "CSV 内容（与 path 二选一）"},
          "path": {"type": "string", "description": "CSV 文件路径"},
          "op": {"type": "string", "description": "table / summary / dedupe，默认 table"}
        },
        "required": []
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "text_stats",
      "description": "文本统计：字数/中文字/字母/数字/句子/非空行",
      "parameters": {
        "type": "object",
        "properties": {
          "text": {"type": "string", "description": "要统计的文本"}
        },
        "required": ["text"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "ip_lookup",
      "description": "查询 IP 归属地（国家/地区/城市/运营商）",
      "parameters": {
        "type": "object",
        "properties": {
          "ip": {"type": "string", "description": "IPv4/IPv6 地址"}
        },
        "required": ["ip"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "mermaid_chart",
      "description": "生成 Mermaid 图代码（flowchart / sequence / pie / gantt）",
      "parameters": {
        "type": "object",
        "properties": {
          "chart_type": {"type": "string", "description": "flowchart / sequence / pie / gantt"},
          "data": {"type": "string", "description": "结构化数据 JSON（见描述）"},
          "title": {"type": "string", "description": "图表标题（可选）"}
        },
        "required": ["chart_type", "data"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "qr_generate",
      "description": "生成二维码（SVG，返回可访问链接）",
      "parameters": {
        "type": "object",
        "properties": {
          "content": {"type": "string", "description": "二维码内容（URL/文本）"}
        },
        "required": ["content"]
      }
    }
  },
  {
    "type": "function",
    "function": {
      "name": "today_hot",
      "description": "获取全网热榜（微博/知乎/百度/抖音/B站/头条）",
      "parameters": {
        "type": "object",
        "properties": {
          "topic": {"type": "string", "description": "weibo/zhihu/baidu/douyin/bilibili/toutiao，默认 weibo"}
        },
        "required": []
      }
    }
  }
]# ── 安全计算（AST 白名单求值，替代 eval，杜绝任意代码执行）──
_MATH_BINOPS = {
    ast.Add: lambda a, b: a + b,
    ast.Sub: lambda a, b: a - b,
    ast.Mult: lambda a, b: a * b,
    ast.Div: lambda a, b: a / b,
    ast.FloorDiv: lambda a, b: a // b,
    ast.Mod: lambda a, b: a % b,
    ast.Pow: lambda a, b: a ** b,
}
_MATH_UNARY = {ast.UAdd: lambda a: +a, ast.USub: lambda a: -a}
_MATH_CONSTANTS = {"pi": math.pi, "e": math.e}
_MATH_FUNCS = {
    "abs": abs, "round": round, "min": min, "max": max, "pow": pow,
    "sqrt": math.sqrt, "sin": math.sin, "cos": math.cos, "tan": math.tan,
    "log": math.log, "log10": math.log10, "log2": math.log2,
}


def safe_calculate(expression: str):
    """只允许数字、四则/幂运算、括号与白名单数学函数/常量的表达式求值"""
    expr = (expression or "").strip()
    if not expr:
        raise ValueError("表达式为空")
    if len(expr) > 200:
        raise ValueError("表达式过长")
    try:
        tree = ast.parse(expr, mode="eval")
    except SyntaxError as e:
        raise ValueError(f"表达式语法错误: {e}") from e

    def eval_node(node):
        if isinstance(node, ast.Expression):
            return eval_node(node.body)
        if isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
            return node.value
        if isinstance(node, ast.BinOp) and type(node.op) in _MATH_BINOPS:
            return _MATH_BINOPS[type(node.op)](
                eval_node(node.left), eval_node(node.right))
        if isinstance(node, ast.UnaryOp) and type(node.op) in _MATH_UNARY:
            return _MATH_UNARY[type(node.op)](eval_node(node.operand))
        if (isinstance(node, ast.Call) and isinstance(node.func, ast.Name)
                and node.func.id in _MATH_FUNCS):
            if len(node.args) > 2:
                raise ValueError("函数参数过多")
            return _MATH_FUNCS[node.func.id](*[eval_node(a) for a in node.args])
        if isinstance(node, ast.Name) and node.id in _MATH_CONSTANTS:
            return _MATH_CONSTANTS[node.id]
        raise ValueError(f"不支持的表达式元素: {type(node).__name__}")

    return eval_node(tree.body)


def is_builtin_tool(name: str) -> bool:
    """判断工具名是否为内置技能"""
    return any(t["function"]["name"] == name for t in TOOLS)


def build_tools(skill_names=None, mcp_servers=None) -> list:
    """根据管理后台下发配置构建当前可用的工具定义列表。

    skill_names 为 None 时保留全部内置技能（向后兼容）；
    传入列表（可为空）时按启用名单过滤，并合并 MCP 工具。
    """
    tools = TOOLS
    if skill_names is not None:
        allowed = set(skill_names)
        tools = [t for t in TOOLS if t["function"]["name"] in allowed]
    if mcp_servers:
        from core.mcp_manager import mcp_manager
        tools = tools + mcp_manager.list_tools(mcp_servers)
    return tools


# ── RSS 源（BBC 在国内被封，改用 DuckDuckGo 新闻搜索兜底）──
_RSS_FEEDS = {
    "tech": [
        "https://www.36kr.com/feed",
        "https://feeds.bbci.co.uk/news/technology/rss.xml",
    ],
    "world": [
        "https://feeds.bbci.co.uk/news/world/rss.xml",
    ],
    "finance": [
        "https://feeds.bbci.co.uk/news/business/rss.xml",
        "https://www.36kr.com/feed",
    ],
    "default": [
        "https://www.36kr.com/feed",
        "https://feeds.bbci.co.uk/news/rss.xml",
    ],
}


# ═════════════════════════════════════════
# 工具执行入口 —— 按名称分发到具体实现
# ═════════════════════════════════════════

def execute_tool(name: str, args: Dict[str, Any]) -> str:
    """执行指定的工具，返回结果字符串"""

    if name == "get_current_time":
        now = datetime.now()
        return f"当前时间: {now.strftime('%Y年%m月%d日 %H:%M:%S')} (星期{['一','二','三','四','五','六','日'][now.weekday()]})"

    if name == "calculate":
        expr = args.get("expression", "")
        try:
            result = safe_calculate(expr)
            return f"计算结果: {expr} = {result}"
        except Exception as e:
            return f"计算失败: {e}"

    if name == "web_search":
        return _duckduckgo_search(args.get("query", ""))

    if name == "url_fetch":
        return _fetch_url(args.get("url", ""))

    if name == "get_weather":
        return _weather(args.get("city", "Beijing"))

    if name == "get_exchange_rate":
        return _exchange_rate(args.get("from_currency", "USD"), args.get("to_currency", "CNY"))

    if name == "wikipedia_lookup":
        return _wiki(args.get("query", ""), args.get("lang", "zh"))

    if name == "date_calc":
        return _date_calc(args.get("date1", ""), args.get("date2"), args.get("add_days"))

    if name == "generate_password":
        length = args.get("length", 16)
        return _gen_password(min(length, 64))

    if name == "make_icon":
        return _make_icon_svg(args.get("description", "icon"), args.get("size", 128))

    if name == "make_chart":
        return _make_chart(args.get("type","bar"), args.get("labels",""), args.get("data",""), args.get("title",""), args.get("unit",""))

    if name == "news_headlines":
        return _news(args.get("topic", "default"))

    if name == "barcode_lookup":
        return _barcode_lookup(args.get("barcode", ""))

    if name == "exchange_convert":
        return _exchange_convert(args.get("from", ""), args.get("to", ""), args.get("amount"))

    if name == "github_search":
        return _github_search(args.get("query", ""))

    if name == "arxiv_search":
        return _arxiv_search(args.get("query", ""))

    if name == "hn_search":
        return _hn_search(args.get("query", ""))

    if name == "pypi_info":
        return _pypi_info(args.get("name", ""))

    if name == "stock_quote":
        return _stock_quote(args.get("symbol", ""))

    if name == "stock_digest":
        return _stock_digest(args.get("symbols", ""))

    if name == "make_table":
        return _make_table(args.get("data", ""), args.get("format", "md"))

    if name == "web_extract":
        return _web_extract(args.get("url", ""), bool(args.get("extract_links", False)))

    if name == "docx_extract":
        return _docx_extract(args.get("path", ""))

    if name == "xlsx_extract":
        return _xlsx_extract(args.get("path", ""), args.get("op", "table"))

    if name == "pptx_extract":
        return _pptx_extract(args.get("path", ""))

    if name == "pdf_extract":
        return _pdf_extract(args.get("path", ""))

    if name == "csv_tools":
        return _csv_tools(args.get("csv_text"), args.get("path"), args.get("op", "table"))

    if name == "text_stats":
        return _text_stats(args.get("text", ""))

    if name == "ip_lookup":
        return _ip_lookup(args.get("ip", ""))

    if name == "mermaid_chart":
        return _mermaid_chart(args.get("chart_type", ""), args.get("data", ""), args.get("title", ""))

    if name == "qr_generate":
        return _qr_generate(args.get("content", ""))

    if name == "today_hot":
        return _today_hot(args.get("topic", "weibo"))

    return f"未知工具: {name}"

# ═══════════════════════════════════════════
# 工具实现（按功能分组）
# ═══════════════════════════════════════════



# --- 网页抓取（P0：域名白名单 + SSRF 防护）---
def _url_fetch_allowed(url: str, allowlist=None) -> tuple:
    """校验抓取目标：协议、域名白名单、内网/保留地址拦截"""
    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        return False, "仅支持 http/https 地址"
    host = (parsed.hostname or "").lower()
    if not host:
        return False, "URL 缺少主机名"
    domains = settings.url_fetch_allowlist if allowlist is None else allowlist
    domains = [d.strip().lower() for d in domains if d.strip()]
    if not domains:
        return False, "抓取工具未启用（未配置域名白名单）"
    if not any(host == d or host.endswith("." + d) for d in domains):
        return False, f"域名不在抓取白名单内: {host}"
    try:
        infos = socket.getaddrinfo(host, None)
    except socket.gaierror:
        return False, f"域名无法解析: {host}"
    for info in infos:
        try:
            ip = ipaddress.ip_address(info[4][0])
        except ValueError:
            continue
        if not ip.is_global:
            return False, f"禁止访问内网/保留地址: {ip}"
    return True, ""


def _fetch_url(url: str) -> str:
    try:
        current = url
        # 手动跟随重定向（最多 3 跳），每一跳都重新校验白名单，防止跳转绕过 SSRF 防护
        for _ in range(4):
            ok, err = _url_fetch_allowed(current)
            if not ok:
                return f"抓取失败: {err}"
            r = requests.get(
                current, headers={"User-Agent": "Mozilla/5.0"},
                timeout=10, allow_redirects=False)
            r.raise_for_status()
            location = r.headers.get("Location")
            if r.is_redirect and location:
                current = urljoin(current, location)
                continue
            # 简单提取正文：去掉 script/style 标签和 HTML
            html = r.text
            html = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
            html = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<[^>]+>', ' ', html)
            text = re.sub(r'\s+', ' ', text).strip()
            return f"【{current}】\n{text[:2000]}"
        return "抓取失败: 重定向次数过多"
    except Exception as e:
        return f"抓取失败: {e}"



# --- 天气 / 汇率 / 百科 ---
def _weather(city: str) -> str:
    try:
        r = requests.get(f"https://wttr.in/{city}?format=%C+%t+%h+%w&lang=zh", timeout=5)
        return f"{city} 天气: {r.text.strip()}"
    except Exception as e:
        return f"天气查询失败: {e}"


def _exchange_rate(frm: str, to: str) -> str:
    try:
        r = requests.get(f"https://api.exchangerate-api.com/v4/latest/{frm.upper()}", timeout=5)
        data = r.json()
        rate = data.get("rates", {}).get(to.upper(), 0)
        return f"1 {frm.upper()} = {rate} {to.upper()}"
    except Exception as e:
        return f"汇率查询失败: {e}"


def _wiki(query: str, lang: str) -> str:
    try:
        r = requests.get(
            f"https://{lang}.wikipedia.org/api/rest_v1/page/summary/{query}",
            headers={"User-Agent": "KnowledgeAgent/1.0"}, timeout=8
        )
        if r.status_code == 200:
            data = r.json()
            return f"【{data.get('title', query)}】\n{data.get('extract', '')[:800]}"
        return f"未找到'{query}'的百科条目"
    except Exception as e:
        return f"百科查询失败: {e}"



# --- 日期计算 / 密码生成 / 新闻 ---
def _date_calc(date1: str, date2: str = None, add_days: int = None) -> str:
    try:
        d1 = datetime.strptime(date1, "%Y-%m-%d")
        if date2:
            d2 = datetime.strptime(date2, "%Y-%m-%d")
            diff = (d2 - d1).days
            return f"{date1} 到 {date2} 间隔 {abs(diff)} 天"
        if add_days is not None:
            result = d1 + timedelta(days=add_days)
            return f"{date1} + {add_days} 天 = {result.strftime('%Y-%m-%d')} (星期{['一','二','三','四','五','六','日'][result.weekday()]})"
        return f"{date1} 是 星期{['一','二','三','四','五','六','日'][d1.weekday()]}"
    except Exception as e:
        return f"日期计算失败: {e}"


def _gen_password(length: int) -> str:
    """生成随机安全密码（大小写字母 + 数字 + 特殊符号）"""
    chars = string.ascii_letters + string.digits + "!@#$%^&*"
    pwd = ''.join(secrets.choice(chars) for _ in range(length))
    return f"随机密码 ({length}位):\n{pwd}"


def _news(topic: str) -> str:
    """获取新闻头条（RSS + DuckDuckGo 兜底）"""
    # 逐个尝试备用 RSS 源（BBC 在部分地区不可达时用 36kr 等兜底）
    feed_urls = _RSS_FEEDS.get(topic, _RSS_FEEDS["default"])
    if isinstance(feed_urls, str):
        feed_urls = [feed_urls]
    for feed_url in feed_urls:
        try:
            r = requests.get(feed_url, headers={"User-Agent": "Mozilla/5.0"}, timeout=5)
            r.raise_for_status()
            items = re.findall(r'<item>.*?<title>(.*?)</title>.*?<description>(.*?)</description>.*?</item>',
                              r.text, re.DOTALL)
            if items:
                topic_names = {"tech": "科技", "world": "国际", "finance": "财经", "default": "综合"}
                lines = [f"📰 最新新闻 ({topic_names.get(topic, topic)}):"]
                for title, desc in items[:8]:
                    title = re.sub(r'<[^>]+>', '', title).strip()
                    lines.append(f"• {title}")
                return "\n".join(lines)
        except Exception:
            continue
    # 🚫 RSS 不可用 → 用 DuckDuckGo 搜索新闻关键词兜底
    try:
        topic_q = {"tech": "科技新闻", "world": "世界新闻", "finance": "财经新闻", "default": "今日新闻"}
        return _duckduckgo_search(topic_q.get(topic, "今日新闻"))
    except Exception:
        return "新闻服务暂时不可用（请检查网络连接）"


def _duckduckgo_search(query: str) -> str:
    """搜索（国内 DDG 不可用，请基于 LLM 已有知识回答）"""
    try:
        resp = requests.get(
            "https://api.duckduckgo.com/",
            params={"q": query, "format": "json", "no_html": 1, "skip_disambig": 1},
            timeout=3
        )
        data = resp.json()
        abstract = data.get("AbstractText", "")
        if abstract:
            return f"搜索结果: {abstract[:800]}"
        related = data.get("RelatedTopics", [])
        if related:
            snippets = [r.get("Text", "")[:200] for r in related[:3] if r.get("Text")]
            return "搜索结果:\n" + "\n".join(f"- {s}" for s in snippets) if snippets else "未找到相关结果"
        return "未找到相关结果"
    except Exception:
        return f"搜索 \"{query[:50]}\" 失败 (国内网络) —— 请基于已有知识回答"


# --- 图表生成（Matplotlib）/ 图标生成（SVG）---
def _make_icon_svg(desc: str, size: int = 128) -> str:
    """根据描述生成简单的 SVG 图标"""
    import os
    desc_lower = desc.lower()
    s = size
    stroke = "#333"
    fill = "#3b82f6"
    if "蓝" in desc or "blue" in desc_lower:
        fill = "#3b82f6"
    elif "绿" in desc or "green" in desc_lower:
        fill = "#22c55e"
    elif "红" in desc or "red" in desc_lower:
        fill = "#ef4444"
    elif "橙" in desc or "orange" in desc_lower:
        fill = "#f59e0b"
    elif "紫" in desc or "purple" in desc_lower:
        fill = "#8b5cf6"

    svg = f"""<?xml version="1.0" encoding="UTF-8"?>
<svg width="{s}" height="{s}" viewBox="0 0 {s} {s}" xmlns="http://www.w3.org/2000/svg">
  <rect width="{s}" height="{s}" rx="{s//8}" fill="{fill}" opacity="0.1"/>
"""
    # 根据关键词选择图形
    if "邮件" in desc or "mail" in desc_lower or "email" in desc_lower:
        svg += f"""  <rect x="{s//4}" y="{s//3}" width="{s//2}" height="{s//3}" rx="{s//32}" fill="none" stroke="{fill}" stroke-width="{s//24}"/>
  <polyline points="{s//4},{s//3} {s//2},{s//2} {s*3//4},{s//3}" fill="none" stroke="{fill}" stroke-width="{s//24}" stroke-linecap="round"/>
"""
    elif "搜索" in desc or "search" in desc_lower or "放大" in desc:
        svg += f"""  <circle cx="{s//2}" cy="{s*2//5}" r="{s//5}" fill="none" stroke="{fill}" stroke-width="{s//20}"/>
  <line x1="{s*3//5}" y1="{s*3//5}" x2="{s*4//5}" y2="{s*4//5}" stroke="{fill}" stroke-width="{s//20}" stroke-linecap="round"/>
"""
    elif "用户" in desc or "user" in desc_lower or "人" in desc or "头像" in desc:
        svg += f"""  <circle cx="{s//2}" cy="{s*2//5}" r="{s//6}" fill="{fill}"/>
  <ellipse cx="{s//2}" cy="{s*4//5}" rx="{s//3}" ry="{s//5}" fill="{fill}"/>
"""
    elif "设置" in desc or "settings" in desc_lower or "齿轮" in desc:
        svg += f"""  <circle cx="{s//2}" cy="{s//2}" r="{s//4}" fill="none" stroke="{fill}" stroke-width="{s//16}"/>
  <circle cx="{s//2}" cy="{s//2}" r="{s//12}" fill="{fill}"/>
"""
    elif "星星" in desc or "star" in desc_lower or "收藏" in desc:
        pts = f"{s//2},{s//10} {s*3//5},{s*2//5} {s*9//10},{s*2//5} {s*2//3},{s*3//5} {s*3//4},{s*9//10} {s//2},{s*7//10} {s//4},{s*9//10} {s//3},{s*3//5} {s//10},{s*2//5} {s*2//5},{s*2//5}"
        svg += f"""  <polygon points="{pts}" fill="{fill}"/>
"""
    elif "主页" in desc or "home" in desc_lower or "房子" in desc:
        svg += f"""  <polyline points="{s//4},{s*3//5} {s//2},{s//5} {s*3//4},{s*3//5}" fill="none" stroke="{fill}" stroke-width="{s//20}" stroke-linecap="round" stroke-linejoin="round"/>
  <rect x="{s*2//5}" y="{s*3//5}" width="{s//5}" height="{s//3}" fill="{fill}"/>
"""
    elif "笔记" in desc or "文档" in desc or "文件" in desc or "doc" in desc_lower:
        svg += f"""  <rect x="{s//4}" y="{s//8}" width="{s//2}" height="{s*3//4}" rx="{s//20}" fill="none" stroke="{fill}" stroke-width="{s//24}"/>
  <line x1="{s*2//5}" y1="{s//2}" x2="{s*3//5}" y2="{s//2}" stroke="{fill}" stroke-width="{s//24}" stroke-linecap="round"/>
  <line x1="{s*2//5}" y1="{s*3//5}" x2="{s*3//5}" y2="{s*3//5}" stroke="{fill}" stroke-width="{s//24}" stroke-linecap="round"/>
"""
    else:
        # 默认：圆形+文字缩写
        svg += f"""  <circle cx="{s//2}" cy="{s//2}" r="{s//3}" fill="{fill}" opacity="0.2"/>
  <text x="{s//2}" y="{s//2}" text-anchor="middle" dy="{s//10}" font-size="{s//3}" fill="{fill}" font-weight="bold">{desc[0]}</text>
"""
    svg += "</svg>"
    # 写文件
    name = "".join(c for c in desc[:20] if c.isalnum() or c in " _-")
    name = name.strip() or "icon"
    out_dir = settings.icon_output_dir or str(
        Path(__file__).resolve().parents[2] / "backend" / "icons")
    os.makedirs(out_dir, exist_ok=True)
    path = os.path.join(out_dir, f"{name}.svg")
    with open(path, "w", encoding="utf-8") as f:
        f.write(svg)
    return f"✅ 图标已生成: {name}.svg\n可通过 http://localhost:8080/icons/{name}.svg 访问\n\nSVG代码:\n```svg\n{svg}\n```"


import json, urllib.parse

def _make_chart(chart_type: str, labels: str, data: str, title: str = "", unit: str = "") -> str:
    """生成图表（基于开源 Matplotlib，本地渲染，无限制）"""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import os, hashlib
    plt.rcParams["font.family"] = "Microsoft YaHei"
    plt.rcParams["axes.unicode_minus"] = False
    labels_list = [l.strip() for l in labels.split(",") if l.strip()]
    data_list = [float(v.strip()) for v in data.split(",") if v.strip()]
    if not labels_list or not data_list:
        return "错误：标签和数据不能为空"
    if len(labels_list) != len(data_list):
        return f"错误：标签数({len(labels_list)})和数据数({len(data_list)})不匹配"
    fig, ax = plt.subplots(figsize=(8, 5))
    colors = ["#3b82f6","#22c55e","#ef4444","#f59e0b","#8b5cf6","#ec4899","#14b8a6","#6366f1"]
    if chart_type == "pie" or chart_type == "doughnut":
        ax.pie(data_list, labels=labels_list, autopct="%1.1f%%", colors=colors[:len(data_list)], startangle=90, wedgeprops={"edgecolor": "white", "linewidth": 2})
        if title: ax.set_title(title, fontsize=14, pad=16)
        if chart_type == "doughnut": fig.gca().add_artist(plt.Circle((0,0),0.58,fc="white"))
    else:
        x = range(len(labels_list))
        if chart_type == "bar" or chart_type == "column":
            bars = ax.bar(x, data_list, color=colors[:len(data_list)], width=0.6, edgecolor="white", linewidth=0.5)
            for bar, v in zip(bars, data_list): ax.text(bar.get_x()+bar.get_width()/2, bar.get_height(), f"{v:.0f}", ha="center", va="bottom", fontsize=10)
        elif chart_type == "barh":
            bars = ax.barh(x, data_list, color=colors[:len(data_list)], height=0.6, edgecolor="white", linewidth=0.5)
            for bar, v in zip(bars, data_list): ax.text(bar.get_width(), bar.get_y()+bar.get_height()/2, f" {v:.0f}", ha="left", va="center", fontsize=10)
            ax.set_yticks(list(x)); ax.set_yticklabels(labels_list, fontsize=11)
            if unit: ax.set_xlabel(unit, fontsize=11)
        elif chart_type == "line":
            ax.plot(x, data_list, "o-", color=colors[0], linewidth=2, markersize=8)
            for i, v in enumerate(data_list): ax.text(i, v+max(data_list)*0.02, f"{v:.0f}", ha="center", fontsize=10)
        else:
            bars = ax.bar(x, data_list, color=colors[:len(data_list)], width=0.6)
            for bar, v in zip(bars, data_list): ax.text(bar.get_x()+bar.get_width()/2, bar.get_height(), f"{v:.0f}", ha="center", va="bottom", fontsize=10)
        if chart_type != "barh":
            ax.set_xticks(list(x)); ax.set_xticklabels(labels_list, fontsize=11)
            if unit: ax.set_ylabel(unit, fontsize=11)
        if title: ax.set_title(title, fontsize=14, pad=12)
        ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
        ax.grid(axis="y", alpha=0.3)
    plt.tight_layout()
    key = hashlib.md5((labels + data + title + chart_type).encode()).hexdigest()[:12]
    # 保存到配置的 charts 目录（默认仓库根 backend/charts，Spring Boot WebConfig 映射）
    import os as _os
    repo_root = Path(__file__).resolve().parents[2]
    chart_dir = settings.chart_output_dir or str(repo_root / "backend" / "charts")
    _os.makedirs(chart_dir, exist_ok=True)
    path = _os.path.join(chart_dir, f"chart_{key}.png")
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)
    base_url = (settings.chart_base_url or "http://localhost:8080").rstrip("/")
    url = f"{base_url}/charts/chart_{key}.png"
    return '<img src="' + url + '" style="max-width:100%;border-radius:8px;margin:8px 0"/><br><small>📊 ' + (title or '图表') + ' · 基于 Matplotlib 渲染</small>'


# ═══════════════════════════════════════════
# 电商 / 数据采集 / 报表 / 汇总 工具（全部免费无 Key，固定可信域名）
# ═══════════════════════════════════════════


def _make_table(data_json: str, format: str = "md") -> str:
    """把 JSON 数组转成 Markdown 表格或 CSV（本地生成，无需网络）"""
    try:
        rows = json.loads(data_json)
    except Exception as e:
        return f"错误：数据不是合法 JSON({e})"
    if not isinstance(rows, list) or not rows:
        return "错误：数据需为非空 JSON 数组"
    first = rows[0]
    if isinstance(first, dict):
        headers = list(first.keys())
        items = [[str(r.get(h, "")) for h in headers] for r in rows]
    elif isinstance(first, list):
        headers = [f"列{i + 1}" for i in range(len(first))]
        items = [[str(c) for c in r] for r in rows]
    else:
        headers = ["值"]
        items = [[str(r)] for r in rows]
    if format == "csv":
        import csv
        import io
        buf = io.StringIO()
        w = csv.writer(buf)
        w.writerow(headers)
        w.writerows(items)
        return buf.getvalue().strip()
    lines = ["| " + " | ".join(headers) + " |",
             "| " + " | ".join("---" for _ in headers) + " |"]
    lines += ["| " + " | ".join(it) + " |" for it in items]
    return "\n".join(lines)


def _validate_symbol(symbol: str) -> bool:
    return bool(symbol) and len(symbol) <= 16 and bool(
        re.fullmatch(r"[A-Za-z0-9^._=:-]+", symbol))


def _fetch_stock(symbol: str) -> dict | None:
    """Stooq 免费 CSV 行情；失败返回 None"""
    if not _validate_symbol(symbol):
        return None
    try:
        r = requests.get(
            f"https://stooq.com/q/l/?s={symbol}&f=sd2t2ohlcv&h&e=csv",
            headers={"User-Agent": "Mozilla/5.0"}, timeout=6)
        r.raise_for_status()
        lines = r.text.strip().splitlines()
        if len(lines) < 2:
            return None
        f = lines[1].split(",")
        if "N/D" in lines[1] or not f:
            return None
        return {
            "symbol": f[0] if f else symbol,
            "date": f[1] if len(f) > 1 else "",
            "open": f[3] if len(f) > 3 else "",
            "high": f[4] if len(f) > 4 else "",
            "low": f[5] if len(f) > 5 else "",
            "close": f[6] if len(f) > 6 else "",
            "volume": f[7] if len(f) > 7 else "",
        }
    except Exception:
        return None


def _stock_quote(symbol: str) -> str:
    d = _fetch_stock(symbol.strip())
    if not d:
        return f"行情获取失败或代码无效: {symbol}"
    return (f"{d['symbol']} {d['date']} 开 {d['open']} 高 {d['high']} "
            f"低 {d['low']} 收 {d['close']} 量 {d['volume']}")


def _stock_digest(symbols: str) -> str:
    codes = [s.strip() for s in symbols.replace("，", ",").split(",") if s.strip()]
    if not codes:
        return "错误：请提供至少一个股票代码"
    rows = [["代码", "日期", "开盘", "最高", "最低", "收盘", "成交量"]]
    failed = []
    for c in codes[:10]:
        d = _fetch_stock(c)
        if not d:
            failed.append(c)
            continue
        rows.append([d["symbol"], d["date"], d["open"], d["high"],
                     d["low"], d["close"], d["volume"]])
    if len(rows) == 1:
        return f"行情获取失败: {', '.join(codes)}"
    md = "| " + " | ".join(rows[0]) + " |\n" + "| " + " | ".join("---" for _ in rows[0]) + " |\n"
    md += "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
    if failed:
        md += f"\n\n以下代码未获取到行情: {', '.join(failed)}"
    return md


def _barcode_lookup(barcode: str) -> str:
    barcode = (barcode or "").strip()
    if not barcode.isdigit() or not (8 <= len(barcode) <= 14):
        return "错误：条码需为 8-14 位数字"
    try:
        r = requests.get(
            f"https://world.openfoodfacts.org/api/v2/product/{barcode}.json",
            headers={"User-Agent": "KnowledgeAgent/1.0"}, timeout=8)
        r.raise_for_status()
        p = r.json().get("product")
        if not p:
            return f"未找到条码 {barcode} 对应的商品信息"
        lines = [f"商品: {p.get('product_name') or p.get('generic_name') or '未知'}"]
        if p.get("brands"):
            lines.append(f"品牌: {p['brands']}")
        if p.get("quantity"):
            lines.append(f"规格: {p['quantity']}")
        if p.get("nutriscore_grade"):
            lines.append(f"营养等级: {p['nutriscore_grade'].upper()}")
        if p.get("ingredients_text"):
            lines.append("配料: " + p["ingredients_text"][:300])
        return "\n".join(lines)
    except Exception as e:
        return f"条码查询失败: {e}"


def _exchange_convert(frm: str, to: str, amount) -> str:
    try:
        amount = float(amount)
    except (TypeError, ValueError):
        return "错误：金额需为数字"
    frm = (frm or "").strip().upper()[:8]
    to = (to or "").strip().upper()[:8]
    if not frm or not to:
        return "错误：请提供源货币与目标货币代码"
    try:
        r = requests.get(
            f"https://api.exchangerate-api.com/v4/latest/{frm}",
            headers={"User-Agent": "KnowledgeAgent/1.0"}, timeout=6)
        r.raise_for_status()
        rate = r.json().get("rates", {}).get(to)
        if rate is None:
            return f"无法获取 {to} 的汇率"
        return f"{amount:g} {frm} = {amount * rate:,.4f} {to} (汇率 {rate})"
    except Exception as e:
        return f"汇率换算失败: {e}"


def _github_search(query: str) -> str:
    q = (query or "").strip()[:80]
    if not q:
        return "错误：搜索关键词不能为空"
    try:
        r = requests.get(
            "https://api.github.com/search/repositories",
            params={"q": q, "per_page": 5, "sort": "stars"},
            headers={"User-Agent": "KnowledgeAgent/1.0",
                     "Accept": "application/vnd.github+json"}, timeout=8)
        r.raise_for_status()
        items = r.json().get("items", [])
        if not items:
            return "未找到相关仓库"
        return "\n".join(
            f"- {it['full_name']} ⭐{it.get('stargazers_count', 0)}: {it.get('description') or '无描述'}"
            for it in items)
    except Exception as e:
        return f"GitHub 搜索失败: {e}"


def _arxiv_search(query: str) -> str:
    q = (query or "").strip()[:100]
    if not q:
        return "错误：搜索关键词不能为空"
    try:
        r = requests.get(
            "http://export.arxiv.org/api/query",
            params={"search_query": f"all:{q}", "max_results": 5,
                    "sortBy": "relevance"},
            headers={"User-Agent": "KnowledgeAgent/1.0"}, timeout=10)
        r.raise_for_status()
        entries = re.findall(r"<entry>(.*?)</entry>", r.text, re.DOTALL)
        if not entries:
            return "未找到相关论文"
        out = []
        for e in entries:
            tm = re.search(r"<title>(.*?)</title>", e, re.DOTALL)
            sm = re.search(r"<summary>(.*?)</summary>", e, re.DOTALL)
            title = re.sub(r"\s+", " ", re.sub(r"<[^>]+>", "", tm.group(1))).strip() if tm else "无标题"
            summary = re.sub(r"\s+", " ", re.sub(r"<[^>]+>", "", sm.group(1))).strip() if sm else ""
            out.append(f"- {title}\n  {summary[:200]}{'…' if len(summary) > 200 else ''}")
        return "\n".join(out)
    except Exception as e:
        return f"arXiv 检索失败: {e}"


def _hn_search(query: str) -> str:
    q = (query or "").strip()[:80]
    if not q:
        return "错误：搜索关键词不能为空"
    try:
        r = requests.get(
            "https://hn.algolia.com/api/v1/search",
            params={"query": q, "tags": "story", "hitsPerPage": 5},
            headers={"User-Agent": "KnowledgeAgent/1.0"}, timeout=8)
        r.raise_for_status()
        hits = r.json().get("hits", [])
        if not hits:
            return "未找到相关讨论"
        return "\n".join(
            f"- {h.get('title', '')} (+{h.get('points', 0)} 分) {h.get('url') or ''}"
            for h in hits)
    except Exception as e:
        return f"Hacker News 检索失败: {e}"


def _pypi_info(name: str) -> str:
    name = (name or "").strip().lower()
    if not name or not re.fullmatch(r"[a-z0-9._-]+", name):
        return "错误：包名不合法"
    try:
        r = requests.get(
            f"https://pypi.org/pypi/{name}/json",
            headers={"User-Agent": "KnowledgeAgent/1.0"}, timeout=8)
        if r.status_code == 404:
            return f"PyPI 上未找到包 {name}"
        r.raise_for_status()
        info = r.json().get("info", {})
        home = info.get("home_page") or (info.get("project_urls") or {}).get("Homepage") or "无"
        return "\n".join([
            f"包: {info.get('name', name)} v{info.get('version', '')}",
            f"描述: {info.get('summary') or '无'}",
            f"主页: {home}",
        ])
    except Exception as e:
        return f"PyPI 查询失败: {e}"


def _web_extract(url: str, extract_links: bool = False, max_chars: int = 6000) -> str:
    """网页正文提取：本地 Trafilatura（Apache-2.0）→ 失败回退 r.jina.ai Reader（免费无 Key）。
    域名需在 KA_CRAWL_ALLOWLIST（未配置时回退 KA_URL_FETCH_ALLOWLIST），内置 SSRF 防护。"""
    url = (url or "").strip()
    allowlist = settings.crawl_allowlist or settings.url_fetch_allowlist
    ok, err = _url_fetch_allowed(url, allowlist)
    if not ok:
        return f"抓取失败: {err}"
    try:
        r = requests.get(url, headers={"User-Agent": "Mozilla/5.0"},
                         timeout=10, allow_redirects=False)
        r.raise_for_status()
        html = r.text
        try:
            import trafilatura
            text = trafilatura.extract(
                html, output_format="markdown",
                include_comments=False, include_tables=True,
                include_links=extract_links)
        except Exception as e:
            logger.info(f"Trafilatura 提取失败: {e}")
            text = None
        if text and text.strip():
            return f"【{url}】\n{text[:max_chars]}"
    except Exception as e:
        logger.info(f"本地抓取失败，尝试 r.jina.ai: {e}")
    # 兜底：r.jina.ai Reader（免费托管，无 Key，带频率限制）
    try:
        rr = requests.get(f"https://r.jina.ai/{url}",
                          headers={"User-Agent": "KnowledgeAgent/1.0"}, timeout=25)
        if rr.status_code == 200 and rr.text.strip():
            return f"【{url}】(via r.jina.ai)\n{rr.text[:max_chars]}"
        return f"网页提取失败: r.jina.ai 返回 {rr.status_code}"
    except Exception as e:
        return f"网页提取失败: {e}"


# ═══════════════════════════════════════════
# 文档提取 / 数据处理 / 工具类技能（本地免费）
# ═══════════════════════════════════════════


def _resolve_allowed_file(path):
    """校验文件路径：必须在 file_access_dirs 内且扩展名合法（防任意文件读取）"""
    import os
    if not path or not str(path).strip():
        return None
    p = Path(str(path).strip()).expanduser()
    if not p.is_absolute():
        p = Path.cwd() / p
    p = p.resolve()
    bases = []
    for d in settings.file_access_dirs:
        b = Path(str(d)).expanduser()
        if not b.is_absolute():
            b = Path.cwd() / b
        bases.append(b.resolve())
    if not any(p == b or b in p.parents for b in bases):
        return None
    if p.suffix.lower() not in (".docx", ".xlsx", ".pptx", ".pdf", ".csv", ".txt", ".md"):
        return None
    return p if p.is_file() else None


def _rows_to_md(rows) -> str:
    """把二维数组转成 Markdown 表格（转义管道符）"""
    if not rows:
        return ""
    esc = lambda c: str(c).replace("|", "\\|").replace("\n", " ")
    lines = ["| " + " | ".join(esc(c) for c in rows[0]) + " |",
             "| " + " | ".join("---" for _ in rows[0]) + " |"]
    for r in rows[1:]:
        lines.append("| " + " | ".join(esc(c) for c in r) + " |")
    return "\n".join(lines)


def _docx_extract(path: str, max_chars: int = 6000) -> str:
    p = _resolve_allowed_file(path)
    if not p:
        return "错误：文件不存在或不在允许目录内"
    from docx import Document
    try:
        doc = Document(str(p))
        parts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        for i, table in enumerate(doc.tables, 1):
            rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
            if rows:
                parts.append(f"表格{i}:\n" + _rows_to_md(rows))
        text = "\n\n".join(parts)
        return text[:max_chars] if text.strip() else "文档中未提取到文本"
    except Exception as e:
        return f"Word 提取失败: {e}"


def _xlsx_extract(path: str, op: str = "table", max_chars: int = 6000) -> str:
    p = _resolve_allowed_file(path)
    if not p:
        return "错误：文件不存在或不在允许目录内"
    import openpyxl
    try:
        wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            rows = [[("" if v is None else str(v)) for v in r]
                    for r in ws.iter_rows(values_only=True)
                    if any(v is not None for v in r)]
            if not rows:
                continue
            if op == "summary":
                header = rows[0]
                numeric = {}
                for r in rows[1:]:
                    for j, v in enumerate(r):
                        if j >= len(header):
                            continue
                        try:
                            val = float(v)
                        except (ValueError, TypeError):
                            continue
                        numeric.setdefault(header[j], []).append(val)
                lines = [f"表: {ws.title}  数据行: {max(len(rows) - 1, 0)}"]
                for h, vals in numeric.items():
                    lines.append(f"- {h}: 合计 {sum(vals):,.2f}  平均 {sum(vals) / len(vals):,.2f} "
                                 f"最大 {max(vals):,.2f}  最小 {min(vals):,.2f}")
                out.append("\n".join(lines))
            else:
                out.append(f"## {ws.title}\n" + _rows_to_md(rows[:50]))
        wb.close()
        return "\n\n".join(out)[:max_chars] or "未读取到数据"
    except Exception as e:
        return f"Excel 读取失败: {e}"


def _pptx_extract(path: str, max_chars: int = 6000) -> str:
    p = _resolve_allowed_file(path)
    if not p:
        return "错误：文件不存在或不在允许目录内"
    from pptx import Presentation
    try:
        prs = Presentation(str(p))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
                elif getattr(shape, "has_table", False):
                    try:
                        rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells]
                                for row in shape.table.rows]
                        if rows:
                            texts.append(_rows_to_md(rows))
                    except Exception:
                        pass
            if texts:
                parts.append(f"## 第{i}页\n" + "\n".join(texts))
        return "\n\n".join(parts)[:max_chars] or "未提取到文本"
    except Exception as e:
        return f"PPT 提取失败: {e}"


def _pdf_extract(path: str, max_pages: int = 20, max_chars: int = 6000) -> str:
    p = _resolve_allowed_file(path)
    if not p:
        return "错误：文件不存在或不在允许目录内"
    try:
        import pymupdf
    except ImportError:
        import fitz as pymupdf
    try:
        doc = pymupdf.open(str(p))
        parts = []
        for i, page in enumerate(doc):
            if i >= max_pages:
                parts.append(f"...（共 {doc.page_count} 页，仅提取前 {max_pages} 页）")
                break
            text = page.get_text().strip()
            if text:
                parts.append(f"## 第{i + 1}页\n{text}")
            try:
                for t in page.find_tables():
                    rows = [[("" if c is None else str(c)) for c in r] for r in t.extract()]
                    if rows:
                        parts.append(_rows_to_md(rows))
            except Exception:
                pass
        doc.close()
        return "\n\n".join(parts)[:max_chars] or "未提取到文本"
    except Exception as e:
        return f"PDF 提取失败: {e}"


def _csv_tools(csv_text=None, path=None, op: str = "table", limit: int = 20) -> str:
    import io
    import csv as _csv
    if csv_text is None and path:
        p = _resolve_allowed_file(path)
        if not p:
            return "错误：文件不存在或不在允许目录内"
        csv_text = p.read_text(encoding="utf-8", errors="replace")
    if not csv_text or not str(csv_text).strip():
        return "错误：CSV 内容为空"
    try:
        rows = list(_csv.reader(io.StringIO(str(csv_text))))
    except Exception as e:
        return f"CSV 解析失败: {e}"
    rows = [[c.strip() for c in r] for r in rows if any(str(c).strip() for c in r)]
    if not rows:
        return "CSV 为空"
    if op == "summary":
        header = rows[0]
        numeric = {}
        for r in rows[1:]:
            for j, v in enumerate(r):
                if j >= len(header):
                    continue
                try:
                    val = float(v)
                except (ValueError, TypeError):
                    continue
                numeric.setdefault(header[j], []).append(val)
        lines = [f"列: {', '.join(header)}  数据行: {len(rows) - 1}"]
        for h, vals in numeric.items():
            lines.append(f"- {h}: 合计 {sum(vals):,.2f}  平均 {sum(vals) / len(vals):,.2f} "
                         f"最大 {max(vals):,.2f}  最小 {min(vals):,.2f}")
        return "\n".join(lines)
    if op == "dedupe":
        seen = set()
        out = [rows[0]]
        for r in rows[1:]:
            k = tuple(r)
            if k not in seen:
                seen.add(k)
                out.append(r)
        return _rows_to_md(out[:limit + 1])
    return _rows_to_md(rows[:limit + 1])


def _text_stats(text: str) -> str:
    text = text or ""
    cn = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
    letters = sum(1 for c in text if c.isalpha())
    digits = sum(1 for c in text if c.isdigit())
    sentences = len(re.findall(r'[^。！？!?；;]+[。！？!?；;]?', text.strip())) if text.strip() else 0
    lines = len([l for l in text.splitlines() if l.strip()])
    return (f"字符 {len(text)} | 中文字 {cn} | 字母 {letters} | 数字 {digits} "
            f"| 句子 {sentences} | 非空行 {lines}")


def _ip_lookup(ip: str) -> str:
    ip = (ip or "").strip()
    if not ip:
        return "错误：请提供 IP 地址"
    try:
        import ipaddress
        ipaddress.ip_address(ip)
    except ValueError:
        return "错误：IP 地址格式不合法"
    try:
        r = requests.get(f"http://ip-api.com/json/{ip}?lang=zh-CN",
                         headers={"User-Agent": "KnowledgeAgent/1.0"}, timeout=6)
        r.raise_for_status()
        d = r.json()
        if d.get("status") != "success":
            return f"查询失败: {d.get('message', '未知')}"
        return (f"{ip} → {d.get('country', '')} {d.get('regionName', '')} {d.get('city', '')}"
                f" | 运营商 {d.get('isp', '')} | {d.get('lat', '')},{d.get('lon', '')}")
    except Exception as e:
        return f"IP 查询失败: {e}"


def _mermaid_chart(chart_type: str, data_json: str, title: str = "") -> str:
    try:
        data = json.loads(data_json) if data_json else {}
    except Exception as e:
        return f"错误：数据不是合法 JSON({e})"
    t = (chart_type or "").lower()
    lines = []
    if t == "pie":
        lines.append("pie title " + (title or "占比"))
        if isinstance(data, dict):
            for k, v in data.items():
                lines.append(f'    "{k}" : {v}')
    elif t in ("flowchart", "flow"):
        lines.append("flowchart TD")
        if isinstance(data, dict):
            for k, v in data.items():
                if k == "edges":
                    continue
                lines.append(f'    {k}["{v}"]')
            for a, b in data.get("edges", []):
                lines.append(f"    {a} --> {b}")
    elif t == "sequence":
        lines.append("sequenceDiagram")
        if isinstance(data, list):
            for item in data:
                if isinstance(item, dict) and item.get("from") and item.get("to"):
                    lines.append(f"    {item['from']}->>{item['to']}: {item.get('text', '')}")
    elif t == "gantt":
        lines.append("gantt")
        lines.append("    dateFormat YYYY-MM-DD")
        if isinstance(data, list):
            for item in data:
                lines.append(f"    {item.get('name', '任务')} : {item.get('start', '')}, {item.get('end', '')}")
    else:
        return "错误：类型仅支持 pie / flowchart / sequence / gantt"
    return "```mermaid\n" + "\n".join(lines) + "\n```"


def _qr_generate(content: str, filename: str | None = None) -> str:
    content = (content or "").strip()
    if not content:
        return "错误：二维码内容不能为空"
    if len(content) > 2000:
        return "错误：内容过长"
    import qrcode
    from qrcode.image.svg import SvgPathImage
    import os as _os
    try:
        qr = qrcode.QRCode(border=2)
        qr.add_data(content)
        qr.make(fit=True)
        img = qr.make_image(image_factory=SvgPathImage)
        name = filename or ("".join(c for c in content[:20] if c.isalnum() or c in " _-").strip() or "qr")
        out_dir = settings.icon_output_dir or str(
            Path(__file__).resolve().parents[2] / "backend" / "icons")
        _os.makedirs(out_dir, exist_ok=True)
        path = _os.path.join(out_dir, f"qr_{name}.svg")
        img.save(path)
        base_url = (settings.chart_base_url or "http://localhost:8080").rstrip("/")
        return f'<img src="{base_url}/icons/qr_{name}.svg" style="max-width:200px;border-radius:8px"/>'
    except Exception as e:
        return f"二维码生成失败: {e}"


def _today_hot(topic: str = "weibo", limit: int = 10) -> str:
    topic = (topic or "weibo").strip().lower()
    mapping = {"weibo": "微博", "zhihu": "知乎", "baidu": "百度", "douyin": "抖音",
               "bilibili": "B站", "bili": "B站", "toutiao": "头条"}
    api_topic = {"weibo": "weibo", "zhihu": "zhihu", "baidu": "baidu", "douyin": "douyin",
                 "bilibili": "bili", "bili": "bili", "toutiao": "toutiao"}.get(topic, "weibo")
    try:
        r = requests.get("https://api.vvhan.com/api/hotlist",
                         params={"type": api_topic},
                         headers={"User-Agent": "Mozilla/5.0"}, timeout=8)
        r.raise_for_status()
        items = (r.json().get("data") or [])[:limit]
        name = mapping.get(topic, "热榜")
        if not items:
            return f"{name}热榜暂无数据"
        lines = [f"🔥 {name}热榜 Top{len(items)}:"]
        for it in items:
            title = it.get("title") or it.get("name") or ""
            hot = it.get("hot") or it.get("hotValue") or ""
            url = it.get("url") or ""
            lines.append(f"- {title} ({hot}) {url}")
        return "\n".join(lines)
    except Exception as e:
        return f"热榜获取失败: {e}"
